import streamlit as st
import pandas as pd
import uuid
import base64
import os
import logging
from io import StringIO, BytesIO
from google.cloud import storage
from google.oauth2 import service_account
from ui.models import AttachmentModel
from streamlit.runtime.uploaded_file_manager import UploadedFile
from ui.services.database_service import SupabaseManager
import email
from email.message import Message
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

class AttachmentComponent:
    def __init__(self):
        self.database_service = SupabaseManager()
        # Initialiser cache for vedleggsinnhold hvis ikke eksisterer
        if 'attachment_cache' not in st.session_state:
            st.session_state.attachment_cache = {}
        # Track which attachment to view
        if 'view_attachment_data' not in st.session_state:
            st.session_state.view_attachment_data = None

    def read_attachment(self, path: str) -> bytes | None:
        """Hent vedleggsinnhold fra cache eller database"""
        cache_key = f"{st.session_state.session_id}_{path}"
        
        # Sjekk cache først
        if cache_key in st.session_state.attachment_cache:
            return st.session_state.attachment_cache[cache_key]
        
        # Hvis ikke i cache, hent fra database
        content_bytes = self.database_service.read_attachment(path=path, bucket_name="attachments")
        
        # Cache resultatet hvis det finnes
        if content_bytes:
            st.session_state.attachment_cache[cache_key] = content_bytes
        
        return content_bytes
    
    def mk_attachment_payload(self, file : UploadedFile, query_id : str) -> AttachmentModel:
        file_id = str(uuid.uuid4())
        ext = file.name.split('.')[-1]
        try:
            content = base64.b64encode(file.getvalue()).decode('utf-8')
        except UnicodeDecodeError as e:
                st.error(f'Error encoding PDF with base64: {e}')
                return None
        except Exception as e:
            st.error(f"Error encoding file {file.name}: {e}")
            return None
        
        # Cache vedleggsinnholdet
        cache_key = f"{st.session_state.session_id}_{file_id}"
        st.session_state.attachment_cache[cache_key] = file.getvalue()

        attachment = AttachmentModel(
            filename=file.name,
            file_id=file_id,
            file_type=file.type,
            path = f'{st.session_state.user_id}/{st.session_state.session_id}/{file_id}.{ext or "dat"}',
            size=file.size,
            content=content,
            query_id=query_id,
        )
        return attachment


    def view_uploaded_file(self, file: UploadedFile):
        """Display an uploaded file - brukes for nylig opplastede filer"""
        file_key = f"{file.name}_{file.size}"
        
        # Cache innholdet ved første tilgang
        if file_key not in st.session_state.attachment_cache:
            st.session_state.attachment_cache[file_key] = file.getvalue()
        
        # Build attachment-like dict with cache_key instead of path
        file_dict = {
            "filename": file.name, 
            "file_type": file.type,
            "size": file.size,
            "cache_key": file_key,  # Use cache_key for uploaded files
        }
        
        st.button(f"📎 {file.name} ({file.size} bytes)", key=f"view_{file_key}", 
                  on_click=lambda f=file_dict: self.on_click_view(f))


    def on_click_view(self, attachment: dict):
        """Set attachment data in session state - rendering happens in render_attachment_dialog"""
        st.session_state.view_attachment_data = attachment

    def render_attachment_dialog(self):
        """Call this once in your main app to render the attachment dialog if one is selected"""
        attachment = st.session_state.get('view_attachment_data')
        if not attachment:
            return

        # Clear immediately so ESC/outside-click doesn't reopen the dialog on next rerun
        st.session_state.view_attachment_data = None

        # Get content - either from cache (uploaded file) or from storage (saved attachment)
        cache_key = attachment.get("cache_key")
        if cache_key:
            content_bytes = st.session_state.attachment_cache.get(cache_key)
        else:
            content_bytes = self.read_attachment(attachment.get("path")) if attachment.get("path") else None

        @st.dialog(f"📎 {attachment.get('filename', 'Vedlegg')}")
        def show_attachment():
            file_type = attachment.get("file_type", "")
            filename = attachment.get("filename", "")

            if file_type == "message/rfc822" and not content_bytes:
                references = attachment.get('reference_paths', '')
                st.info(f'This view shows in total {len(references)} uploaded' + (' emails' if len(references) > 1 else ' email')) if references else st.info('No references found for this email attachment.')
                #if st.button("show email", )

                # Render email from DB data directly
                st.markdown(f"**From:** {attachment.get('from_addr', '')}")
                st.markdown(f"**To:** {', '.join(attachment.get('to', [])) if isinstance(attachment.get('to'), list) else attachment.get('to', '')}")
                st.markdown(f"**Subject:** {attachment.get('subject', '')}")
                st.markdown(f"**Date:** {attachment.get('date', '')}")
                
                st.divider()
                st.text(attachment.get("body", ""))
            elif content_bytes:
                if file_type == "application/pdf":
                    RenderDocBytes.parse_pdf(content_bytes, metadata={"filename": filename})
                elif file_type in ["text/plain", "text/markdown"]:
                    RenderDocBytes.parse_txt(content_bytes, metadata={"filename": filename})
                elif file_type == "text/csv":
                    RenderDocBytes.parse_csv(content_bytes, metadata={"filename": filename})
                elif file_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
                    RenderDocBytes.parse_excel(content_bytes, metadata={"filename": filename})
                elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    RenderDocBytes.parse_docx(content_bytes, metadata={"filename": filename})
                elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    RenderDocBytes.parse_pptx(content_bytes, metadata={"filename": filename})
                elif file_type == "message/rfc822":
                    RenderDocBytes.parse_email(content_bytes, metadata={"filename": filename})
                else:
                    RenderDocBytes.parse_generic(content_bytes, metadata={"filename": filename})
            else:
                st.error(f'Kunne ikke hente vedlegg: {filename}')

            if st.button("Lukk"):
                st.rerun()

        show_attachment()

    def view_attachment(self, attachment: dict, content_bytes: bytes | None = None, key: str = None, sig_icon: str = None):
        """Display attachments from session history"""
        if not attachment:
            logger.warning("view_attachment called with None attachment")
            return

        file_id = attachment.get('file_id', "unknown_id")
        filename = attachment.get('filename', 'Ukjent fil')
        file_date = attachment.get('file_date', '') or attachment.get("date")
        #st.json(attachment)

        # Use stable key based on file_id, not random uuid
        button_key = key if key else f"att_{file_id}"

        label = f"{sig_icon} 📅 {file_date} 📎 {filename}" if sig_icon else f"📎 {filename}"
        st.button(label, key=button_key, on_click=lambda a=attachment: self.on_click_view(a))


class RenderDocBytes:
    """Hjelpeklasse for å parse bytes til forskjellige formater"""
    @staticmethod
    def parse_csv(content: bytes, metadata: dict = None) -> pd.DataFrame | None:
        try:
            df = pd.read_csv(StringIO(content.decode('utf-8', errors='ignore')))
            st.dataframe(df)
        except Exception as e:
            logger.error(f'Error parsing CSV: {e}', exc_info=True)
            st.error(f'Error parsing CSV: {e}')

    @staticmethod
    def parse_excel(content: bytes, metadata: dict = None) -> pd.DataFrame | None:
        try:
            df = pd.read_excel(BytesIO(content))
            st.dataframe(df)
        except Exception as e:
            logger.error(f'Error parsing Excel: {e}', exc_info=True)
            st.error(f'Error parsing Excel: {e}')
        
    @staticmethod
    def parse_pdf(content: bytes, metadata: dict = None) -> bytes | None:
        try:
            st.pdf(BytesIO(content))
        except Exception as e:
            logger.error(f'Error displaying PDF: {e}', exc_info=True)
            st.error(f'Error displaying PDF: {e}')
    
    @staticmethod
    def parse_txt(content: bytes, metadata: dict = None) -> str | None:
        try:
            st.text(content.decode('utf-8', errors='ignore'))
        except Exception as e:
            logger.error(f'Error parsing text file: {e}', exc_info=True)
            st.error(f'Error parsing text file: {e}')
    
    @staticmethod
    def parse_docx(content: bytes, metadata: dict = None) -> str | None:
        try:
            document = Document(BytesIO(content))
        except Exception as e:
            logger.error(f'Error loading DOCX file: {e}', exc_info=True)
            st.error(f'Error loading DOCX file: {e}')
            return
        
        try:
            full_text = []
            for para in document.paragraphs:
                full_text.append(para.text)
            st.text("\n".join(full_text))
        except ImportError:
            st.error("python-docx library is required to parse .docx files. Please install it with `pip install python-docx`.")
        except Exception as e:
            logger.error(f'Error parsing DOCX file: {e}', exc_info=True)
            st.error(f'Error parsing DOCX file: {e}')

    @staticmethod
    def parse_pptx(content: bytes, metadata: dict = None) -> str | None:
        try:
            presentation = Presentation(BytesIO(content))
        except Exception as e:
            logger.error(f'Error loading PPTX file: {e}', exc_info=True)
            st.error(f'Error loading PPTX file: {e}')
            return
        
        try:
            full_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            st.text("\n".join(full_text))
        except ImportError:
            st.error("python-pptx library is required to parse .pptx files. Please install it with `pip install python-pptx`.")
        except Exception as e:
            logger.error(f'Error parsing PPTX file: {e}', exc_info=True)
            st.error(f'Error parsing PPTX file: {e}')
    
    @staticmethod
    def parse_email(content: bytes, metadata: dict = None) -> str | None:
        try:
            msg = email.message_from_bytes(content)
            st.markdown(f"**From:** {msg.get('From')}")
            st.markdown(f"**To:** {msg.get('To')}")
            st.markdown(f"**Subject:** {msg.get('Subject')}")
            st.divider()
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        st.text(part.get_payload(decode=True).decode('utf-8', errors='ignore'))
            else:
                st.text(msg.get_payload(decode=True).decode('utf-8', errors='ignore'))
        except Exception as e:
            logger.error(f'Error parsing email file: {e}', exc_info=True)
            st.error(f'Error parsing email file: {e}')


    @staticmethod
    def parse_generic(content: bytes, metadata: dict = None) -> str | None:
        try:
            st.text(content.decode('utf-8', errors='ignore'))
        except Exception as e:
            logger.error(f'Error parsing file: {e}', exc_info=True)
            st.error(f'Error parsing file: {e}')

    
def get_attachment_component() -> AttachmentComponent:
    """Get AttachmentComponent instance (not cached to ensure code updates work)"""
    return AttachmentComponent()