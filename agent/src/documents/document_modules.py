import os
from io import BytesIO
import tempfile
from PyPDF2 import PdfReader
import logging
import base64
from typing import List, Optional
from datetime import datetime
import uuid

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from models import VectorStoreMetadata,FileType, AttachmentModel, EmailModel, WriteEmail, WriteDocx
import ocrmypdf
from email.message import Message
import email
from docx import Document as DocxDocument
from pptx import Presentation

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class BaseHandler:
    """Handles ONLY parsing - no storage logic."""
    
    def __init__(self, chunk_size: int = 800, chunk_overlap: int = 100):
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )

       
    
    @staticmethod
    def to_plain_text(docs: List[Document]) -> str:
        """Extract concatenated text from documents."""
        return "\n\n".join([d.page_content for d in docs])
    
    @staticmethod
    def to_dict(docs: List[Document]) -> List[dict]:
        """Convert to dict for JSON/BigQuery."""
        return [{"content": d.page_content, "metadata": d.metadata} for d in docs]



class EmailHandler(BaseHandler):
    def __init__(self):
        super().__init__()

    def _dedoce_base64(self, content: str) -> bytes:
        try:
            return base64.b64decode(content)
        except Exception as e:
            logger.error(f"Error decoding base64 content: {e}", exc_info=True)
            raise ValueError("Invalid base64 content") from e

    def _extract_email_body(self, msg : Message) -> dict:
        if msg.is_multipart():
            body_text = ""
            body_html = None
            for part in msg.walk():
                content_type = part.get_content_type()
                if content_type == "text/plain":
                    body_text += part.get_payload(decode=True).decode(part.get_content_charset() or "utf-8")
                elif content_type == "text/html":
                    body_html = part.get_payload(decode=True).decode(part.get_content_charset() or "utf-8")
        else:
            body_text = msg.get_payload(decode=True).decode(msg.get_content_charset() or "utf-8")
            body_html = None
        return {"html" : body_html, "text": body_text}

    def _extract_attachments(self, msg : Message) -> list:
        allowed_types = ["application/pdf", 
                         "application/msword", 
                         "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                         "text/plain",
                         "text/csv",
                         "message/rfc822",

                         ]
        attachments = []
        for part in msg.walk():
            content_disposition = part.get("Content-Disposition")
            if content_disposition and "attachment" in content_disposition:
                filename = part.get_filename()
                if filename:
                    if part.get_content_type() in allowed_types:
                        payload = part.get_payload(decode=True)
                        try:
                            content = payload.decode(part.get_content_charset() or "utf-8")
                        except (UnicodeDecodeError, LookupError):
                            content = base64.b64encode(payload).decode("ascii")
                        attachments.append({
                            "filename": filename,
                            "file_type": part.get_content_type(),
                            "size" : len(payload),
                            "file_id": str(uuid.uuid4()),
                            "content": content
                        })
                    else:
                        logger.warning(f"Attachment '{filename}' has unsupported content type '{part.get_content_type()}', skipping. Allowed types are: {allowed_types}")
                else:
                    logger.warning("Attachment part found without filename, skipping.")
        return attachments
 
    def _extract_email_data(self, msg : Message, file_id : str, query_id : str, user_id: str , session_id : str) -> dict:
        #file_id = str(uuid.uuid4())
        attachments_list = self._extract_attachments(msg)
        attachments = []
        att_ids = []
        if attachments_list:
            for att in attachments_list:
                ext = os.path.splitext(att["filename"])[1].lower()
                attachment_model = AttachmentModel(
                    filename=att["filename"],
                    file_id=att["file_id"],
                    file_type=att["file_type"],
                    size=att["size"],
                    content=att["content"],
                    query_id=query_id,
                    event_id=None,
                    path = f'{user_id}/{session_id}/{att.get("file_id")}{ext}',
                )
                attachments.append(attachment_model)
                att_ids.append(att["file_id"])
        body = self._extract_email_body(msg)
        refs = msg.get("References")
        email_size = len(msg.as_bytes()) if hasattr(msg, 'as_bytes') else len(msg.as_string().encode('utf-8'))
        email_data = EmailModel(
                file_id=file_id,
                path = f'{user_id}/{session_id}/{file_id}.eml',
                query_id=query_id,

                subject=msg.get("Subject", ""),
                from_addr=msg.get("From", ""),
                to=[addr.strip() for addr in msg.get("To", "").split(",")],
                cc=[addr.strip() for addr in msg.get("Cc", "").split(",")] if msg.get("Cc") else None,
                bcc=[addr.strip() for addr in msg.get("Bcc", "").split(",")] if msg.get("Bcc") else None,
                date=email.utils.parsedate_to_datetime(msg.get("Date")) if msg.get("Date") else None,

                message_id=msg.get("Message-ID"),
                in_reply_to=msg.get("In-Reply-To"),
                references=refs,
                thread_id=msg.get("Thread-ID"),
                thread_index=msg.get("Thread-Index"),
                thread_topic=msg.get("Thread-Topic"),

                body_text=body.get("text", ""),
                body_html=body.get("html"),
                headers=dict(msg.items()) if msg.items() else None,
                size=email_size,

                attachments=att_ids if att_ids else None,
            )

        return {"email" : email_data, "attachments" : attachments if attachments else []}

    def parse_eml_to_obj(self, content: str, user_id, query_id, session_id, file_id) -> dict:
        '''Process EML content and extract email data and attachments
        
        Args:
            content (str): The EML content as a base64 encoded string.
            user_id (str): The ID of the user associated with the email.
            query_id (str): The ID of the query associated with the email.
            session_id (str): The ID of the session associated with the email.
            file_id (str): The ID of the file associated with the email.
        Returns:
            dict: A dictionary containing the extracted email data and attachments.
        
        '''
        content = self._dedoce_base64(content)
        try:
            msg = email.message_from_bytes(content)
        except Exception as e:
            logger.error(f"Error parsing EML content: {e}", exc_info=True)
            raise ValueError("Invalid EML content") from e
        email_data = self._extract_email_data(msg, 
                                              query_id=query_id, 
                                              user_id=user_id, 
                                              session_id=session_id, 
                                              file_id=file_id)
        return email_data
    
    def parse_eml_to_docs(self, content: bytes, metadata : dict) -> list[Document]:
        '''Process EML content and extract email data and attachments
        
        Args:
            content (str): The EML content as a base64 encoded string.
            user_id (str): The ID of the user associated with the email.
            query_id (str): The ID of the query associated with the email.
            session_id (str): The ID of the session associated with the email.
        Returns:
            dict: A dictionary containing the extracted email data and attachments.
        
        '''
        try:
            msg = email.message_from_bytes(content)
        except Exception as e:
            logger.error(f"Error parsing EML content: {e}", exc_info=True)
            raise ValueError("Invalid EML content") from e
        body = self._extract_email_body(msg)
        chunks = self.splitter.split_text(body.get("text", ""))
        docs = []
        if not chunks:
            logger.warning("No text chunks created from email body.")
            chunks = [body.get("text", "")]
        metadata_all = {**metadata,
                        "file_size": len(content),
                        "file_type": "message/rfc822",
                        "creator" : msg.get("From"),
                        "title" : msg.get("Subject"),
                        "created_at" : email.utils.parsedate_to_datetime(msg.get("Date")) if msg.get("Date") else None,
                        } 
        metadata_model = VectorStoreMetadata.model_validate(metadata_all)
        
        for i, chunk in enumerate(chunks):
            #logger.debug(f"Chunk {i + 1}/{len(chunks)}: {chunk[:100]}...")  # Log first 100 chars of each chunk
            docs.append(Document(page_content=chunk, metadata={**metadata_model.model_dump(), "chunk": i+1, "total_chunks": len(chunks)}))
        return docs
    
    def mk_eml(self, email_data : WriteEmail) -> bytes:
        msg = email.message.EmailMessage()
        msg["Subject"] = email_data.subject
        msg["From"] = email_data.from_addr
        msg["To"] = ", ".join(email_data.to)
        if email_data.cc:
            msg["Cc"] = ", ".join(email_data.cc)
        if email_data.bcc:
            msg["Bcc"] = ", ".join(email_data.bcc)
        msg.set_content(email_data.body)
        return msg.as_bytes()

class PDFHandler(BaseHandler):
    def __init__(self):
        super().__init__()

    
    def _safe_pdf_date(self, metadata, field: str) -> Optional[datetime]:
        """Access a PyPDF2 metadata date property safely, returning None on parse errors."""
        try:
            return getattr(metadata, field)
        except Exception:
            return None

    def _needs_ocr(self, content: bytes) -> bool:
        """Detect if PDF needs OCR based on text density."""
        try:
            reader = PdfReader(BytesIO(content))
            total_pages = len(reader.pages)
            pages_with_text = 0
            total_text_length = 0
            
            for page in reader.pages:
                text = page.extract_text().strip()
                if text and len(text) > 50:  # More than metadata/page numbers
                    pages_with_text += 1
                    total_text_length += len(text)
            
            # Heuristikk: Hvis < 50% av sidene har tekst ELLER lite meningsfullt innhold per side
            # 500 chars/page threshold skiller ekte innhold fra bare metadata/sidenumre/overskrifter
            text_coverage = pages_with_text / total_pages if total_pages > 0 else 0
            avg_text_per_page = total_text_length / total_pages if total_pages > 0 else 0

            return text_coverage < 0.5 or avg_text_per_page < 500
        except Exception as e:
            logger.warning(f"Could not analyze PDF for OCR need: {e}")
            return False  # Default to no OCR if detection fails

    def _ocr_bytes(self, content: bytes) -> bytes:
        """OCR a PDF with improved error handling."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as inp, \
            tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as out:
            
            try:
                inp.write(content)
                inp.flush()
                inp.close()  # Close before OCR reads it
                
                ocrmypdf.ocr(
                    inp.name,
                    out.name,
                    deskew=True,  # Deskew pages for better OCR accuracy. Not compatibel with redo_ocr
                    #redo_ocr=True,  # Re-OCR entire document for better text extraction. Not compatibel with deskew
                    skip_text=False,  # Keep existing text
                    optimize=1,  # Light optimization
                    force_ocr=False  # Don't re-OCR text pages
                )
                
                with open(out.name, 'rb') as f:
                    return f.read()
                    
            except ocrmypdf.exceptions.PriorOcrFoundError:
                logger.info("PDF already has OCR text, returning original")
                return content
            except Exception as e:
                logger.error(f"OCR failed: {e}, returning original PDF")
                return content
            finally:
                # Cleanup temp files
                for f in [inp.name, out.name]:
                    try:
                        os.unlink(f)
                    except:
                        pass
    
    def parse_pdf_to_docs(self, content: bytes, metadata: dict) -> List[Document]:
        if self._needs_ocr(content):
            logger.info("PDF needs OCR, processing...")
            content = self._ocr_bytes(content)
        
        count_without_text = 0
        try:
            reader = PdfReader(BytesIO(content))
        except Exception as e:
            logger.error(f"Error reading PDF: {e}")
            return []
        docs = []
        base_meta = metadata | {
            "creator": reader.metadata.creator if reader.metadata else None,
            "producer": reader.metadata.producer if reader.metadata else None,
            "created_at": self._safe_pdf_date(reader.metadata, "creation_date") if reader.metadata else None,
            "updated_at": self._safe_pdf_date(reader.metadata, "modification_date") if reader.metadata else None,
            "title": reader.metadata.subject or reader.metadata.title if reader.metadata else None,
            "keywords": reader.metadata.get("/Keywords") if reader.metadata else None,
            "file_size": len(content),
            "file_type": "application/pdf",
        }
        meta_model = VectorStoreMetadata.model_validate(base_meta)

        for i, page in enumerate(reader.pages):
            txt = page.extract_text().strip() if page.extract_text() else ""
            if not txt:
                count_without_text += 1
                logger.debug(f"Page {i + 1} has no extractable text.")
                continue

            else:
                docs.append(Document(
                    page_content=txt,
                    metadata={**meta_model.model_dump(), "chunk": i + 1, "total_chunks": len(reader.pages)}
                ))
            
        if not docs or count_without_text == len(reader.pages):
            logger.warning("No pages extracted from PDF.")
            return []
        logger.debug(f"Extracted {len(docs)} pages with text out of {len(reader.pages)} total pages.")
        return docs
    
class TextHandler(BaseHandler):
    def __init__(self):
        super().__init__()

    def parse_text_to_docs(self, content : bytes, metadata: dict) -> List[Document]:
        text = content.decode('utf-8', errors='ignore')
        try:
            chunks = self.splitter.split_text(text)
        except Exception as e:
            logger.error(f"Error splitting text: {e}")
            chunks = [text]  # Fallback to whole text if splitting fails

        if not chunks:
            logger.warning("No chunks created from text.")
            return []
        
        metadata_all = {**metadata,
                        "file_size": len(content),
                        "file_type": "text/plain",
                        }
        meta_model = VectorStoreMetadata.model_validate(metadata_all)
        try:
            return [
                Document(
                    page_content=chunk,
                    metadata={**meta_model.model_dump(), "chunk": i+1, "total_chunks": len(chunks)}
                )
                for i, chunk in enumerate(chunks)
            ]
        except Exception as e:
            logger.error(f"Error creating Document objects: {e}")
            return []
        

    def parse_csv_to_docs(self, content: bytes, metadata: dict) -> list[Document]:
        content_decoded = content.decode('utf-8', errors='ignore')
        docs = []
        chunks = self.splitter.split_text(content_decoded)
        if not chunks:
            logger.warning("No chunks created from CSV content.")
            chunks = [content_decoded]

        metadata_all = {**metadata,
                        "file_size": len(content),
                        "file_type": "text/csv",
                        }
        metadata_model = VectorStoreMetadata.model_validate(metadata_all)

        for i, chunk in enumerate(chunks):
            #logger.debug(f"Chunk {i + 1}/{len(chunks)}: {chunk[:100]}...")  # Log first 100 chars of each chunk
            docs.append(Document(page_content=chunk, metadata={**metadata_model.model_dump(), "chunk": i+1, "total_chunks": len(chunks)}))
        return docs
        
class DocxHandler(BaseHandler):
    def __init__(self):
        super().__init__()

    def parse_docx_to_docs(self, content: bytes, metadata: dict) -> list[Document]:
        docs = []
        try:
            file = BytesIO(content)
            word_doc = DocxDocument(file)
        except Exception as e:
            logger.error(f"Error loading bytes content to python-docx Document: {e}")
            return []
        
        props = word_doc.core_properties
        metadata_full = {
                    **metadata,
                    "title": props.title,
                    "creator": props.author,
                    "created_at": props.created.isoformat() if props.created else None,         
                    "updated_at": props.modified.isoformat() if props.modified else None,
                    "comments": props.comments,
                    "language": props.language,
                    "file_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",

                    #"subject": props.subject,
                    #"keywords": props.keywords,
                    #"revision": props.revision,
                    #"last_modified_by": props.last_modified_by,
                    #"category": props.category,
                    #"content_status": props.content_status,
                }
        metadata_model = VectorStoreMetadata.model_validate(metadata_full)
        
        for i, para in enumerate(word_doc.paragraphs):
                text = para.text.strip()
                if text:
                    docs.append(Document(
                        page_content=text,
                        metadata={**metadata_model.model_dump(), "chunk": i+1, "total_chunks": len(word_doc.paragraphs)}))
                    
        return docs

    def mk_docx(self, docx_data : WriteDocx) -> bytes:
        doc = DocxDocument()
        if docx_data.heading:
            doc.add_heading(docx_data.heading, level=1)
        for para in docx_data.paragraphs:
            doc.add_paragraph(para)
        with BytesIO() as output:
            doc.save(output)
            return output.getvalue()
        
class XlsxHandler(BaseHandler):
    def __init__(self):
        super().__init__()
    
    def parse_xlsx_to_docs(self, content: bytes, metadata: dict) -> list[Document]:
        logger.warning("XLSX parsing not implemented yet.")
        return []
    
class PptxHandler(BaseHandler):
    def __init__(self):
        super().__init__()

    def parse_pptx_to_docs(self, content: bytes, metadata: dict) -> list[Document]:
        docs = []
        try:
            file = BytesIO(content)
            ppt_doc = Presentation(file)
        except Exception as e:
            logger.error(f"Error loading bytes content to python-pptx Presentation: {e}")
            return []
        props = ppt_doc.core_properties
        metadata_full = {
                    **metadata,
                    "title": props.title,
                    "creator": props.author,
                    "created_at": props.created.isoformat() if props.created else None,         
                    "updated_at": props.modified.isoformat() if props.modified else None,
                    "comments": props.comments,
                    "language": props.language,
                    "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",}
        metadata_model = VectorStoreMetadata.model_validate(metadata_full)
        for i, slide in enumerate(ppt_doc.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text.strip())
            slide_text_combined = "\n".join(slide_text).strip()
            if slide_text_combined:
                docs.append(Document(
                    page_content=slide_text_combined,
                    metadata={**metadata_model.model_dump(), "chunk": i+1, "total_chunks": len(ppt_doc.slides)}))
        return docs


class DocumentProcessor(BaseHandler):
    """Handles ONLY parsing - no storage logic."""
    
    def __init__(self,):
        pass

    def parse(self, content : str, 
                           file_type : FileType,
                           metadata : dict = None) -> List[Document]:
        '''Main entry point for parsing attachments. Handles different file types and routes to appropriate parsers.
        Args:
            content (str): The content of the attachment, expected to be a base64 encoded string.
            file_type (str): The MIME type of the file, used to determine parsing method.
            metadata (dict): Additional metadata to attach to each Document, such as file_id and session_id.
        Returns:
            List[Document]: A list of Document objects extracted from the attachment, ready for embedding and storage.
        '''

        if "file_id" not in metadata or "session_id" not in metadata or "embedding_model" not in metadata:
            logger.error("Metadata must include 'file_id', 'session_id', and 'embedding_model'")
            raise ValueError("Metadata must include 'file_id', 'session_id', and 'embedding_model'")
        
        try:
            content_decoded = base64.b64decode(content)
        except Exception as e:
            logger.error(f"Failed to decode base64 content for attachment {metadata.get('file_id')}: {e}")
            return []
        
        if file_type == "application/pdf":
            return PDFHandler().parse_pdf_to_docs(content_decoded, metadata=metadata)
        elif file_type in ["text/plain", "text/markdown"]:
            return TextHandler().parse_text_to_docs(content_decoded, metadata=metadata)
        elif file_type == "text/csv":
            return TextHandler().parse_csv_to_docs(content_decoded, metadata=metadata)
        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return XlsxHandler().parse_xlsx_to_docs(content_decoded, metadata=metadata)
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return DocxHandler().parse_docx_to_docs(content_decoded, metadata=metadata)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return PptxHandler().parse_pptx_to_docs(content_decoded, metadata=metadata)
        elif file_type == "message/rfc822":
            return EmailHandler().parse_eml_to_docs(content_decoded, metadata=metadata)
        else:
            logger.warning(f"Unsupported file type {file_type} for attachment {metadata.get('file_id')}")
            return []




    
